

from utils import *

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

import io

import matplotlib.pyplot as plt


BAR_COLORS = ["#826fc2","#001f80","#4d9b1e","#f865c6","#ecd378","#ba004c","#8f4400","#f65656"]*10000

def create_pptx_v4(df):
    # Creates a PowerPoint presentation with:
    # Numerical columns displayed using a histogram, boxplot, and table.
    # Categorical columns displayed using a bar chart and table.
    #
    # PowerPoint settings
    ppt_file = "EDA_Presentation_v4.pptx"
    prs = Presentation()

    width = 13.33
    slide_width = Inches(width)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Loop through all columns in the DataFrame
    
    for count, col_name in enumerate(df.columns):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        title = slide.shapes.title
        title.text = f"{col_name}"
        
        if title and title.text_frame:
            for paragraph in title.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT

        col_type = df[col_name].dtype

        # Numerical column
        if col_type in ['int64', 'float64']:
            fig = box_hist_with_table(df, col_name, BAR_COLORS[count])  # Generate the combined plot for numerical data

        # Categorical column
        elif col_type == 'object':
            fig = create_categorical_plot_with_table(df, col_name, BAR_COLORS[count])  # Generate the combined plot for categorical data

        else:
            continue  # Skip columns that are neither numeric nor categorical

        # Save plot to a buffer
        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format="png", dpi=300)
        img_buffer.seek(0)
        plt.close(fig)

        # Add the image to the slide
        left = Inches(0.0)
        top = Inches(1.5)

        slide.shapes.add_picture(img_buffer, left, top, width=Inches(width))  # Adjust width to fit the slide

    # Save PowerPoint to a buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer.getvalue()


def box_hist_with_table(df,col_name,BAR_COLOR, width = 13.33, height = 6  ):
    # Creates a combined plot with:
    #
    #   A box plot on the top-left
    #   A histogram on the bottom-left
    #   A summary table on the right
    
    import matplotlib.pyplot as plt
    import matplotlib.gridspec as gridspec

    # Calculate statistics for the table
    stats = calculate_numeric_stats(df, col_name)  # Use your utility function here
    stats.columns = ["Statistic", "Value"]  # Ensure proper column labels

    data = df[col_name].dropna()

    # Create the figure
    fig = plt.figure(figsize=(width, height))  # Adjust the figure size as needed
    gs = gridspec.GridSpec(2, 3, width_ratios=[2, 1, 2], height_ratios=[1, 3])  # Adjust grid layout

    # Box Plot (Top-Left)
    ax_box = fig.add_subplot(gs[0, 0:2])
    ax_box.grid(axis='x',zorder=1.0)
    ax_box.boxplot(data, vert=False, patch_artist=True, boxprops=dict(facecolor=BAR_COLOR), zorder=2.0)

    ax_box.axes.get_yaxis().set_visible(False)  # Hide y-axis for box plot
    ax_box.set_xticklabels([])                  # Remove x-axis labels
    ax_box.spines[['top','left', 'right', 'bottom']].set_visible(False)
    ax_box.set_title(f"Histogram of {col_name}", fontsize = 15)

    # Histogram (Bottom-Left)
    ax_hist = fig.add_subplot(gs[1, 0:2])
    ax_hist.grid(axis='x' ,zorder=3.0)
    ax_hist.hist(data, bins=20, color=BAR_COLOR, edgecolor='black', zorder = 4.0)
    
    ax_hist.set_xlabel(col_name, fontsize=13)
    ax_hist.set_ylabel("Count", fontsize = 13)
    for label in (ax_hist.get_xticklabels() + ax_hist.get_yticklabels()): label.set_fontsize(11)
    ax_hist.spines[['top','right']].set_visible(False)
    

    # Table (Right)
    ax_table = fig.add_subplot(gs[:, 2])  # Span rows 0 and 1 for the table
    ax_table.axis("off")  # Turn off the axis
    table = ax_table.table(
        cellText=stats.values,
        colLabels=stats.columns,
        loc="center",
        cellLoc="center")
    
    table.auto_set_font_size(False)
    table.set_fontsize(14)
    table.auto_set_column_width(col=list(range(len(stats.columns))))

    table.scale(1 , 1.75)

    # Adjust layout to avoid overlap
    plt.tight_layout()

    return fig


def create_categorical_plot_with_table(df, col_name, BAR_COLOR, width=13.33, height=6):
    """
    Creates a combined plot with:
    - A bar chart on the left
    - A summary table of the top 6 most frequent values on the right
    """
    import matplotlib.pyplot as plt
    import matplotlib.gridspec as gridspec

    # Calculate bar chart data using your utility function
    bar_data = bar_chart_data(df, col_name, top_n_rows=6)

    # Truncate y-axis labels
    max_label_length = 12  # Maximum length for y-axis labels
    bar_data[col_name] = bar_data[col_name].apply(lambda x: x[:max_label_length] + "..." if len(x) > max_label_length else x)

    # Create the figure 
    fig = plt.figure(figsize=(width, height))  # Adjust figure size as needed
    gs = gridspec.GridSpec(1, 2, width_ratios=[1.1, 1])  # Grid layout for bar chart and table

    # Bar Chart (Left)
    ax_bar = fig.add_subplot(gs[0, 0])
    ax_bar.grid(axis='x', zorder=1.0)
    ax_bar.barh(bar_data[col_name], bar_data["Occurrences"], color=BAR_COLOR, edgecolor="black", zorder=2.0)
    ax_bar.set_title(f"Bar Chart of {col_name}", fontsize=15)
    ax_bar.set_xlabel("Occurrences", fontsize=13)
    ax_bar.set_ylabel(col_name, fontsize=13)
    for label in (ax_bar.get_xticklabels() + ax_bar.get_yticklabels()):
        label.set_fontsize(11)

    ax_bar.invert_yaxis()  # Invert y-axis for better readability
    ax_bar.spines[['top', 'right']].set_visible(False)

    # Table (Right)
    ax_table = fig.add_subplot(gs[0, 1])
    ax_table.axis("off")  # Turn off axis for table
    table = ax_table.table(
        cellText=bar_data.values,
        colLabels=bar_data.columns,
        loc="center",
        cellLoc="center"
    )

    # Adjust table properties
    table.auto_set_font_size(False)
    table.set_fontsize(12)

    # Dynamically adjust column widths
    col_width_scale = max(1, width / 13.33 * 1.2)  # Scale columns based on figure width
    table.scale(col_width_scale, 1.75)

    # Truncate long text in the first column
    for cell in table.get_celld().values():
        cell_text = cell.get_text().get_text()
        if len(cell_text) > 20:  # Limit text to 20 characters
            truncated_text = cell_text[:11] + "..."
            cell.get_text().set_text(truncated_text)

    # Adjust layout to avoid overlap
    plt.tight_layout()

    return fig
