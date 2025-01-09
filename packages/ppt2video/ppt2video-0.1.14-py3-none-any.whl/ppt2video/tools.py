from dataclasses import dataclass, field
from pptx import Presentation
from google.cloud import texttospeech_v1beta1 as tts
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
import os
import re
import win32com.client
from moviepy.editor import CompositeVideoClip, VideoFileClip
from PIL import Image 
Image.ANTIALIAS=Image.LANCZOS

@dataclass
class Meta:
    # PPT settings
    ppt_file: str # name does not need to include .pptx 
    ppt_path: str = 'data/ppt/'  # Directory for the PPT and image files
    image_prefix: str = 'slide'  # The prefix for image file names (used when saving slides as images)
    image_extension: str = 'png'  # The image file format 
    ppt_extension: str = '.pptx'  # The PowerPoint file extension
    convert_slides_upto_slide_no: int = 0   # Convert to video only slide number upto this
    save_slide_images: bool = True # Save slides automatically from PPT, although there are already saved slides (works only under windows, if you have PowerPoint software)

    # Google TTS settings
    voice_enabled: bool = True  # Enable or disable voice narration
    google_application_credentials: str = None  # Location of the Google API key (downloaded as JSON)
    voice_path: str = 'data/voice/'  # Directory to save the generated audio files
    max_size: int = 4500  # Maximum text size limit for a single Google TTS API request (default 5000)
    slide_break: float = 1.0  # Time delay (in seconds) between slides
    line_break: float = 0.5  # Time delay (in seconds) when there's a line break in the text (e.g., '\n')
    lang: str = 'E'  # Language setting: 'E' for English, 'K' for Korean 
    wave: bool = True  # Whether to use Wavenet voices (True or False)
    wave_E: str = 'D'
    wave_K: str = 'C'
    speaking_rate_EN: float = 1.2 # English 
    speaking_rate_KR: float = 1.2 # Korean

    # MoviePy video settings
    fps: int = 24  # Frames per second for the video
    fade_duration: float = 0.15 # Slide fade duration
    fade_after_slide: list = field(default_factory=list) # fade effect after given slide number: starting from 0
    target_slide_for_video: list = field(default_factory=list)
    video_file_path: list = field(default_factory=list)
    video_height_scale: list = field(default_factory=list)
    video_location: list = field(default_factory=list)
    video_interrupt: bool = False

def ppt_to_video(meta: Meta): 
    if not os.path.exists(meta.ppt_path):
        os.makedirs(meta.ppt_path)
    
    meta.ppt_file = (meta.ppt_file if meta.ppt_file.endswith(meta.ppt_extension) else meta.ppt_file + meta.ppt_extension)

    if meta.save_slide_images:
        save_ppt_as_images(meta)

    if meta.voice_enabled:
        if meta.google_application_credentials == None:
            print('*****')
            print('Need to set up Google Cloud Authentication')
            print('Please refer to the README.md')
            print('*****')
            return None

        if not os.path.exists(meta.voice_path):
            os.makedirs(meta.voice_path)
        num = ppt_to_text(meta)
        timepoints, total_duration = ppt_tts(meta, num)
        # video_from_ppt_and_voice(meta, timepoints)
        composite_video_from_ppt_and_voice(meta, timepoints)
    else:
        num = ppt_to_text(meta)
        video_from_ppt(meta, num)

def _clean_text(input_text):
    # Ensure UTF-8 compatibility: decode and encode to handle encoding correctly
    input_text = input_text.encode('utf-8').decode('utf-8')

    # Remove non-Korean, non-English chars, non-numbers, and special characters 
    # except commas, periods, question marks, exclamation marks, spaces, %, $, &, -
    input_text = re.sub(r'[^a-zA-Z0-9가-힣.,?!%\$\n\s&-/]', '', input_text)

    # Remove all newlines
    input_text = input_text.replace('\n', '') 

    # Replace multiple spaces with a single space
    input_text = re.sub(r'\s+', ' ', input_text)

    # Return cleaned text
    return input_text.strip()

def _write_to_file(content, current_file_number, current_size, meta: Meta):
    txt_file = f"{os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, ''))}_{current_file_number}.txt"

    mode = 'w' if current_size == 0 else 'a'
    with open(txt_file, mode, encoding='utf-8') as notes_file:
        notes_file.write(content)
    
    return current_size + len(content.encode('utf-8'))

def ppt_to_text(meta: Meta):
    ppt = Presentation(os.path.join(meta.ppt_path, meta.ppt_file))
    if not meta.voice_enabled:
        return len(ppt.slides)

    header = '''<speak>\n'''
    footer = '''</speak>'''
    file_number = 0
    current_size = _write_to_file(header, file_number, 0, meta)

    mark_separator = '.' # for the Google TTS English engine, MARK tag seems need to be followed by a char.
    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            notes = _clean_text(notes)
            slide_content = f'<mark name="slide{slide_number}"/>{mark_separator}\n<break time="{round(meta.slide_break/2,1)}s"/>\n'
            slide_content += notes.replace('\n', f'\n<break time="{meta.line_break}s"/>\n') + f'\n<break time="{meta.slide_break}s"/>\n'
        else:
            slide_content = f'<mark name="slide{slide_number}"/>{mark_separator}\n<break time="{meta.slide_break}s"/>\n'

        if current_size + len(slide_content.encode('utf-8')) > meta.max_size:
            _write_to_file(footer, file_number, current_size, meta)
            file_number += 1
            current_size = 0
            slide_content = header + slide_content

        current_size = _write_to_file(slide_content, file_number, current_size, meta)
        if meta.convert_slides_upto_slide_no > 0 and slide_number == meta.convert_slides_upto_slide_no:
            break

    _write_to_file(footer, file_number, current_size, meta)
    txt_file_number = file_number+1

    return txt_file_number 

def get_text_script_path(meta: Meta, n: int):
    return f"{os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, '_'+str(n)+'.txt'))}"

def get_voice_file_path(meta: Meta, n: int):
    return os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, '_'+str(n)+'.wav'))

def ppt_tts(meta: Meta, txt_file_number: int):
    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = meta.google_application_credentials

    client = tts.TextToSpeechClient()
    if meta.lang == 'E':
        language_code = 'en-US' 
        speaking_rate = meta.speaking_rate_EN
        name = 'en-US-Wavenet-' + meta.wave_E
    elif meta.lang == 'K':
        language_code = 'ko-KR' 
        speaking_rate = meta.speaking_rate_KR
        name = 'ko-KR-Wavenet-' + meta.wave_K
    else: # default
        language_code = 'en-US' 
        speaking_rate = meta.speaking_rate_EN
        name = 'en-US-Wavenet-' + meta.wave_E
    
    if meta.wave == True: # WaveNet voice (1 mil words/month vs 4 mil in basic)
        voice = tts.VoiceSelectionParams(language_code=language_code, name=name, ssml_gender=tts.SsmlVoiceGender.MALE)
    else:
        voice = tts.VoiceSelectionParams(language_code=language_code, ssml_gender=tts.SsmlVoiceGender.MALE)

    audio_config = tts.AudioConfig(
        audio_encoding=tts.AudioEncoding.LINEAR16,
        speaking_rate = speaking_rate
    )
    
    timepoint_dict = {}
    total_duration = 0
    for i in range(txt_file_number):
        txt_file = get_text_script_path(meta, i)
        voice_file = get_voice_file_path(meta, i)

        with open(txt_file, 'r', encoding='utf-8') as file:
            text_content = file.read()

        synthesis_input = tts.SynthesisInput(ssml=text_content)
        request = tts.SynthesizeSpeechRequest(
            input=synthesis_input,
            voice=voice,
            audio_config=audio_config, 
            enable_time_pointing=[tts.SynthesizeSpeechRequest.TimepointType.SSML_MARK]
        )
        response = client.synthesize_speech(request=request)

        with open(voice_file, "wb") as out:
            out.write(response.audio_content)
            print(voice_file + ' done')

        with AudioFileClip(voice_file) as voice_clip:
            total_duration += voice_clip.duration

        timepoint_list = []
        if response.timepoints:
            for time_point in response.timepoints:
                print(f'Mark name: {time_point.mark_name}, Time: {time_point.time_seconds} seconds')
                timepoint_list.append([int(time_point.mark_name[5:]), time_point.time_seconds]) # mark_name = 'slide#'
        else:
            print('No timepoints found.')
        timepoint_dict[voice_file] = timepoint_list

    print(f"Total duration of voice files is {total_duration}")

    return timepoint_dict, total_duration


def video_from_ppt_and_voice(meta: Meta, timepoints, fps=24):
    images_path = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension,''))
    output_file = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '.mp4'))
    videos_with_diff_audio_files = []

    for audio_file, slide_times in timepoints.items():
        audio_clip = AudioFileClip(audio_file)

        video_clips = []
        for i in range(len(slide_times)):
            start_time = slide_times[i][1]  # Get the start time for the slide
            if i < len(slide_times)-1:
                end_time = slide_times[i + 1][1]  # Get the end time for the next slide
            else:
                end_time = audio_clip.duration
            slide_number = slide_times[i][0]

            # Construct the image filename
            slide_image_filename = f'{meta.image_prefix}{slide_number}.{meta.image_extension}'
            slide_image_path = os.path.join(images_path, slide_image_filename)

            # Load the slide image
            slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)

            # Apply fade-out to the current slide if it's in fade_after_slide list
            fade_after_slide_next_one = [i+1 for i in meta.fade_after_slide]

            if slide_number in meta.fade_after_slide:
                slide_clip = slide_clip.fadeout(meta.fade_duration)
            if slide_number in fade_after_slide_next_one: 
                slide_clip = slide_clip.fadein(meta.fade_duration)

            video_clips.append(slide_clip)

        # Concatenate video clips for the current audio
        video_for_an_audio_file = concatenate_videoclips(video_clips)
        video_for_an_audio_file = video_for_an_audio_file.set_audio(audio_clip)
        videos_with_diff_audio_files.append(video_for_an_audio_file)

    # Concatenate all videos into one final video
    final_video = concatenate_videoclips(videos_with_diff_audio_files)

    # Set fps for the final video
    final_video.fps = fps
    
    # final_video.write_videofile(output_file, codec="libx264")
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video with audio generated and saved')

def composite_video_from_ppt_and_voice(meta: Meta, timepoints, fps=24):
    images_path = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension,''))
    output_file = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '.mp4'))
    videos_with_diff_audio_files = []

    for audio_file, slide_times in timepoints.items():
        audio_clip = AudioFileClip(audio_file)

        video_clips = []
        for i in range(len(slide_times)):
            start_time = slide_times[i][1] # Get the start time for the slide
            if i < len(slide_times)-1:
                end_time = slide_times[i + 1][1] # Get the end time for the next slide
            else:
                end_time = audio_clip.duration 
            slide_number = slide_times[i][0]

            # Construct the image filename
            slide_image_filename = f'{meta.image_prefix}{slide_number}.{meta.image_extension}'
            slide_image_path = os.path.join(images_path, slide_image_filename)

            # Load the slide image

            if slide_number in meta.target_slide_for_video:
                slide_clip = ImageClip(slide_image_path)
                
                ith = meta.target_slide_for_video.index(slide_number)
                vc_path = os.path.join(meta.ppt_path, meta.video_file_path[ith])
                video_overlay = VideoFileClip(vc_path)
                lv = video_overlay.duration
                if lv + 0.5 > end_time-start_time: 
                    raise Exception(f'Composite Video Duration Error at slide: {slide_number}')
                video_overlay = video_overlay.set_duration(end_time-start_time)

                print(f'---------')
                print(f'pricessing slide {slide_number}')
                print(f'slide size (w, h) = ({slide_clip.w}, {slide_clip.h})')
                print(f'original video size (w, h) = ({video_overlay.w}, {video_overlay.h})')
                if meta.video_interrupt == True: 
                    user_input = input("Type 'y' to continue or any other key to halt the process: ")
                    if user_input.lower() == 'y':
                        print("Continuing the process...")
                    else:
                        raise Exception(f'Halting the process...')
                print(f'---------')

                video_overlay = video_overlay.resize(height=slide_clip.h*meta.video_height_scale[ith])  
                video_overlay = video_overlay.set_position(meta.video_location[ith])  
                
                # Composite the video on top of the slide image
                slide_clip = CompositeVideoClip([slide_clip, video_overlay])
                slide_clip = slide_clip.set_duration(end_time - start_time).set_start(start_time)
            else: 
                slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)

            # Apply fade-out to the current slide if it's in fade_after_slide list
            fade_after_slide_next_one = [i+1 for i in meta.fade_after_slide]

            if slide_number in meta.fade_after_slide:
                slide_clip = slide_clip.fadeout(meta.fade_duration)
            if slide_number in fade_after_slide_next_one: 
                slide_clip = slide_clip.fadein(meta.fade_duration)

            video_clips.append(slide_clip)

        # Concatenate video clips for the current audio
        video_for_an_audio_file = concatenate_videoclips(video_clips)
        video_for_an_audio_file = video_for_an_audio_file.set_audio(audio_clip)
        videos_with_diff_audio_files.append(video_for_an_audio_file)

    # Concatenate all videos into one final video
    final_video = concatenate_videoclips(videos_with_diff_audio_files)

    # Set fps for the final video
    final_video.fps = fps
    
    # final_video.write_videofile(output_file, codec="libx264")
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video with audio generated and saved')

def video_from_ppt(meta: Meta, num_slides: int, fps=24):
    images_path = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension,''))
    output_file = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '.mp4'))

    video_clips = []
    for i in range(num_slides):
        start_time = i*meta.slide_break
        end_time = start_time+meta.slide_break

        slide_image_filename = f'{meta.image_prefix}{i}.{meta.image_extension}'
        slide_image_path = os.path.join(images_path, slide_image_filename)

        slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)
        video_clips.append(slide_clip)

    final_video = concatenate_videoclips(video_clips)
    final_video.fps = fps
    
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video generated and saved')


def save_ppt_as_images(meta: Meta):
    slide_folder = os.path.abspath(os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '')))
    os.makedirs(slide_folder, exist_ok=True)

    # Initialize PowerPoint
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_file = os.path.abspath(os.path.join(meta.ppt_path, meta.ppt_file))
    presentation = ppt_app.Presentations.Open(ppt_file, WithWindow=False)

    # Loop through slides and save each as an image
    for i, slide in enumerate(presentation.Slides):
        image_path = os.path.join(slide_folder, f'{meta.image_prefix}{i}.{meta.image_extension}')
        slide.Export(image_path, "PNG")
        print(f"Saved slide {i} to {image_path}")

    # Close the presentation
    presentation.Close()
    ppt_app.Quit()
