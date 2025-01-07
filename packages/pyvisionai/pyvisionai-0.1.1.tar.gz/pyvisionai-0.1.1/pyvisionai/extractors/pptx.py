"""PPTX content extractor."""

import os
from pptx import Presentation
from PIL import Image
import io

from pyvisionai.extractors.base import BaseExtractor


class PptxTextImageExtractor(BaseExtractor):
    """Extract text and images from PPTX files."""

    def save_image(self, image_data: bytes, output_dir: str, image_name: str) -> str:
        """Save an image to the output directory."""
        try:
            # Convert image data to PIL Image
            image = Image.open(io.BytesIO(image_data))
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            # Save as JPEG (supported format)
            img_path = os.path.join(output_dir, f"{image_name}.jpg")
            image.save(img_path, "JPEG", quality=95)
            return img_path
        except Exception as e:
            print(f"Error saving image: {str(e)}")
            raise

    def extract(self, pptx_path: str, output_dir: str) -> str:
        """Process PPTX file by extracting text and images."""
        try:
            pptx_filename = os.path.splitext(os.path.basename(pptx_path))[0]
            prs = Presentation(pptx_path)

            # Generate markdown content
            md_content = f"# {pptx_filename}\n\n"

            # Process slides
            for slide_num, slide in enumerate(prs.slides, 1):
                md_content += f"## Slide {slide_num}\n\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        md_content += f"{shape.text}\n\n"

                # Extract images
                image_count = 0
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_count += 1
                            image_name = f"{pptx_filename}_slide_{slide_num}_image_{image_count}"
                            img_path = self.save_image(shape.image.blob, output_dir, image_name)
                            
                            # Get image description using configured model
                            image_description = self.describe_image(img_path)
                            md_content += f"[Image {image_count}]\n"
                            md_content += f"Description: {image_description}\n\n"
                            
                            # Clean up image file
                            os.remove(img_path)
                        except Exception as e:
                            print(f"Error processing image in slide {slide_num}: {str(e)}")
                            continue

            # Save markdown file
            md_file_path = os.path.join(output_dir, f"{pptx_filename}.md")
            with open(md_file_path, "w", encoding="utf-8") as md_file:
                md_file.write(md_content)

            return md_file_path

        except Exception as e:
            print(f"Error processing PPTX: {str(e)}")
            raise 